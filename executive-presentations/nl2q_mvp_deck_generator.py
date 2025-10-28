"""
NL2Q Analyst MVP Deck Generator - Conexus Corporate Template Edition
Showcase intelligent visualization capabilities with Conexus corporate design

Author: NL2Q Analytics Team
Date: October 10, 2025
Version: 2.0 - Corporate Template with Modern Design
"""

import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Optional


class NL2QMVPDeckGenerator:
    """
    Generate MVP showcase presentation using Conexus Corporate Template
    """
    
    def __init__(self, template_path: str = r"assets\Conexus Corporate Template 2025.pptx",
                 assets_dir: str = "assets", output_dir: str = "outputs"):
        self.template_path = Path(template_path)
        self.assets_dir = Path(assets_dir)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"❌ Conexus template not found: {self.template_path}")
        
        # Modern professional color palette - Azure/Conexus inspired
        self.colors = {
            "azure_blue": RGBColor(0, 120, 212),        # Primary Azure
            "azure_light": RGBColor(80, 180, 240),      # Light Azure
            "azure_dark": RGBColor(0, 90, 158),         # Dark Azure
            "primary": RGBColor(59, 130, 246),          # Blue 500
            "primary_light": RGBColor(96, 165, 250),    # Blue 400
            "accent_purple": RGBColor(139, 92, 246),    # Purple 500
            "accent_cyan": RGBColor(6, 182, 212),       # Cyan 600
            "success": RGBColor(16, 185, 129),          # Emerald 500
            "warning": RGBColor(251, 146, 60),          # Orange 400
            "white": RGBColor(255, 255, 255),
            "gray_50": RGBColor(249, 250, 251),
            "gray_100": RGBColor(243, 244, 246),
            "gray_200": RGBColor(229, 231, 235),
            "gray_300": RGBColor(209, 213, 219),
            "gray_700": RGBColor(55, 65, 81),
            "text_dark": RGBColor(30, 41, 59),
            "text_medium": RGBColor(71, 85, 105),
            "text_light": RGBColor(148, 163, 184),
        }
        
        # Conexus template layout indices
        self.layouts = {
            'title': 0,           # Title slide
            'blue_bar': 1,        # Content with blue bar
            'two_column': 6,      # Two column layout
            'section_blue': 14,   # Section divider (blue)
            'section_orange': 15, # Section divider (orange)
            'closing': 16         # Closing slide
        }
        
        # Screenshot filenames
        self.screenshots = [
            "Analyst-1.png",
            "Analyst-2.png", 
            "Analys-3.png",
            "Analyst-4.png"
        ]
        
    def create_presentation(self) -> Presentation:
        """Load Conexus Corporate Template (preserves backgrounds, logos, disclaimers)"""
        print(f"📋 Loading Conexus Corporate Template from: {self.template_path}")
        prs = Presentation(str(self.template_path))
        
        # Remove example slides but keep master slides and layouts intact
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        print(f"✅ Template loaded with {len(prs.slide_layouts)} slide layouts")
        return prs
    
    def add_title_slide(self, prs: Presentation):
        """Create professional title slide using Conexus template"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['title']])
        
        # Find and update title and subtitle placeholders
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Main title placeholder
                if 'title' in shape.name.lower() and 'placeholder 21' not in shape.name.lower():
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = "NL2Q Analyst"
                    p.font.size = Pt(60)
                    p.font.bold = True
                    p.font.color.rgb = self.colors['azure_blue']
                    p.alignment = PP_ALIGN.LEFT
                
                # Subtitle placeholder (placeholder 21)
                elif 'placeholder 21' in shape.name.lower():
                    shape.text_frame.clear()
                    
                    # Line 1: Main subtitle
                    p = shape.text_frame.paragraphs[0]
                    p.text = "Intelligent Visualization & Natural Language Query Platform"
                    p.font.size = Pt(24)
                    p.font.color.rgb = self.colors['text_dark']
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(12)
                    
                    # Line 2: MVP tag
                    p2 = shape.text_frame.add_paragraph()
                    p2.text = "MVP Showcase • October 2025"
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = self.colors['azure_light']
                    p2.alignment = PP_ALIGN.LEFT
        
        return slide
        
    def add_section_divider(self, prs: Presentation, title: str, use_orange: bool = False):
        """Add section divider using Conexus template"""
        layout_key = 'section_orange' if use_orange else 'section_blue'
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts[layout_key]])
        
        # Update title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = title
                break
        
        return slide
    
    def add_problem_statement_slide(self, prs: Presentation):
        """Slide 2: Problem Statement"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['blue_bar']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = "The Challenge"
                break
        
        # Problem description
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.2)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        p = content_frame.paragraphs[0]
        p.text = "Business users struggle with complex data visualization requirements:"
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['text_dark']
        p.font.bold = True
        p.space_after = Pt(20)
        
        # Pain points with icons
        pain_points = [
            ("📊", "Static dashboards can't adapt to natural language queries"),
            ("🔍", "Comparing categories requires manual filtering and analysis"),
            ("⏰", "Temporal patterns are missed without intelligent context detection"),
            ("🎨", "Generic visualizations fail to highlight key insights"),
        ]
        
        y_position = 3.2
        for icon, point in pain_points:
            self._add_bullet_point(slide, icon, point, Inches(1.8), Inches(y_position))
            y_position += 0.9
        
    def add_solution_slide(self, prs: Presentation):
        """Slide 3: Solution Overview"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['blue_bar']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = "The NL2Q Solution"
                break
        
        # Solution description
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.0)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        p = content_frame.paragraphs[0]
        p.text = "AI-powered intelligent visualization with adaptive layouts and contextual insights"
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['text_dark']
        p.font.bold = True
        p.space_after = Pt(20)
        
        # Key features in boxes
        features = [
            ("🤖 LLM-Driven Planning", "GPT-4 analyzes queries to generate optimal visualization plans"),
            ("🎯 Smart Comparisons", "Automatically detects comparison queries and creates filtered KPIs"),
            ("📈 Contextual Insights", "Temporal patterns and category breakdowns with priority ranking"),
            ("⚡ Dynamic Filtering", "Real-time data filtering based on user interactions"),
        ]
        
        x_positions = [1.5, 7.0]
        y_positions = [3.2, 5.0]
        
        for idx, (title, desc) in enumerate(features):
            x = x_positions[idx % 2]
            y = y_positions[idx // 2]
            self._add_feature_box(slide, title, desc, Inches(x), Inches(y), Inches(5.0), Inches(1.5))
    
    def add_architecture_slide(self, prs: Presentation):
        """Slide 4: Technical Architecture"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['blue_bar']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = "Technical Architecture"
                break
        
        # Architecture layers
        layers = [
            {
                "title": "🎨 Frontend Layer",
                "tech": "React + TypeScript + Plotly.js",
                "components": ["AdaptiveLayout", "ComparisonCards", "KPICards", "Dynamic Charts"]
            },
            {
                "title": "🧠 Intelligence Layer", 
                "tech": "GPT-4o-mini + AsyncOpenAI",
                "components": ["VisualizationPlanner", "Temporal Detection", "Filter Generation"]
            },
            {
                "title": "⚙️ Backend Layer",
                "tech": "Python + FastAPI + Pandas",
                "components": ["Dynamic Orchestrator", "Query Engine", "Schema Tools"]
            },
            {
                "title": "💾 Data Layer",
                "tech": "Azure SQL + Vector Store",
                "components": ["DWHPRODIBSA", "Schema Cache", "Query History"]
            }
        ]
        
        y_position = 2.0
        for layer in layers:
            self._add_architecture_layer(slide, layer, Inches(1.5), Inches(y_position), Inches(10.333), Inches(1.1))
            y_position += 1.25
            
    def add_screenshot_slide(self, prs: Presentation, screenshot_path: Path, title: str, description: str):
        """Add slide with screenshot and description"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['blue_bar']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = title
                break
        
        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8), Inches(10.333), Inches(0.6)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = description
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors['text_dark']
        
        # Screenshot (if exists)
        if screenshot_path.exists():
            try:
                # Add screenshot with border
                pic = slide.shapes.add_picture(
                    str(screenshot_path),
                    Inches(1.5), Inches(2.6),
                    width=Inches(10.333)
                )
                
                # Add subtle shadow/border effect
                pic.shadow.inherit = False
            except Exception as e:
                print(f"⚠️ Could not add screenshot {screenshot_path}: {e}")
                # Add placeholder text
                placeholder = slide.shapes.add_textbox(
                    Inches(4.0), Inches(4.0), Inches(5.333), Inches(1.0)
                )
                placeholder_frame = placeholder.text_frame
                p = placeholder_frame.paragraphs[0]
                p.text = f"Screenshot: {screenshot_path.name}"
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors['text_light']
    
    def add_key_features_slide(self, prs: Presentation):
        """Slide: Key Features Showcase"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['blue_bar']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = "Key Features Delivered"
                break
        
        features = [
            {
                "icon": "✅",
                "title": "Dynamic Comparison KPIs",
                "desc": "Auto-generates filtered KPIs for each comparison category (YES/NO, Regions, etc.)",
                "status": "Delivered"
            },
            {
                "icon": "✅",
                "title": "Temporal Context Detection",
                "desc": "95% confidence detection of 7 temporal patterns (YoY, QoQ, MoM, WoW, etc.)",
                "status": "Delivered"
            },
            {
                "icon": "✅",
                "title": "Contextual Insights Sidebar",
                "desc": "Priority-ranked comparison cards with click-to-filter interactions",
                "status": "Delivered"
            },
            {
                "icon": "✅",
                "title": "Professional UI Design",
                "desc": "Enterprise-grade blue color scheme with optimized spacing and typography",
                "status": "Delivered"
            },
            {
                "icon": "⏳",
                "title": "Full Drilldown Filtering",
                "desc": "Complete chart filtering on comparison card clicks (logging implemented)",
                "status": "In Progress"
            },
            {
                "icon": "📋",
                "title": "Export & Sharing",
                "desc": "PDF/PowerPoint export of visualizations and insights",
                "status": "Planned"
            }
        ]
        
        y_position = 2.0
        x_positions = [1.5, 7.0]
        
        for idx, feature in enumerate(features):
            x = x_positions[idx % 2]
            y = y_position + (idx // 2) * 1.5
            self._add_feature_card(slide, feature, Inches(x), Inches(y), Inches(5.0), Inches(1.3))
    
    def add_demo_flow_slide(self, prs: Presentation):
        """Slide: Demo Flow"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['blue_bar']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = "Demo Flow"
                break
        
        # Flow steps
        steps = [
            ("1️⃣", "User Query", "Compare prescribers with PDRP flag enabled vs disabled"),
            ("2️⃣", "AI Processing", "GPT-4 analyzes query, detects comparison pattern"),
            ("3️⃣", "Plan Generation", "Creates 6 filtered KPIs (TRX/NRX/LunchLearn × YES/NO)"),
            ("4️⃣", "Visualization", "Renders adaptive layout with charts and contextual insights"),
            ("5️⃣", "Interaction", "User clicks comparison cards to filter and explore data"),
        ]
        
        y_position = 2.0
        for icon, title, desc in steps:
            self._add_flow_step(slide, icon, title, desc, Inches(2.0), Inches(y_position), Inches(9.333))
            y_position += 1.0
    
    def add_results_slide(self, prs: Presentation):
        """Slide: Results & Impact"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['blue_bar']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = "Results & Impact"
                break
        
        # Metrics in cards
        metrics = [
            ("95%", "Detection Accuracy", "Temporal pattern recognition"),
            ("100%", "Dynamic Filtering", "KPIs correctly filtered by category"),
            ("7", "Temporal Patterns", "YoY, QoQ, MoM, WoW, DoD, etc."),
            ("6", "Components Built", "React TypeScript components"),
        ]
        
        x_positions = [1.5, 4.5, 7.5, 10.5]
        y_position = 2.5
        
        for idx, (value, title, desc) in enumerate(metrics):
            self._add_metric_card(slide, value, title, desc, Inches(x_positions[idx]), Inches(y_position), Inches(2.5), Inches(2.0))
        
        # Benefits
        benefits_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(5.0), Inches(10.333), Inches(2.0)
        )
        benefits_frame = benefits_box.text_frame
        benefits_frame.word_wrap = True
        
        p = benefits_frame.paragraphs[0]
        p.text = "Business Impact:"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.colors['text_dark']
        p.space_after = Pt(12)
        
        benefits = [
            "⚡ Instant insights from natural language queries without manual dashboard configuration",
            "🎯 Accurate comparisons with automatic category detection and filtering",
            "📊 Professional, enterprise-grade visualizations that adapt to query context",
            "🚀 Scalable to any comparison dimension (regions, products, time periods, etc.)",
        ]
        
        for benefit in benefits:
            p = benefits_frame.add_paragraph()
            p.text = benefit
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text_dark']
            p.space_after = Pt(8)
            p.level = 0
    
    def add_next_steps_slide(self, prs: Presentation):
        """Slide: Next Steps & Roadmap"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['two_column']])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = "Next Steps & Roadmap"
                break
        
        # Immediate next steps
        immediate_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.0), Inches(5.0), Inches(4.5)
        )
        immediate_frame = immediate_box.text_frame
        immediate_frame.word_wrap = True
        
        p = immediate_frame.paragraphs[0]
        p.text = "🎯 Immediate (Sprint 1-2)"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.space_after = Pt(12)
        
        immediate_items = [
            "Complete drilldown filtering on card clicks",
            "Add export to PDF/PowerPoint",
            "Implement chart type switching",
            "Add data table toggle view",
        ]
        
        for item in immediate_items:
            p = immediate_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text_dark']
            p.space_after = Pt(8)
        
        # Short-term roadmap
        shortterm_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(2.0), Inches(5.0), Inches(4.5)
        )
        shortterm_frame = shortterm_box.text_frame
        shortterm_frame.word_wrap = True
        
        p = shortterm_frame.paragraphs[0]
        p.text = "🚀 Short-term (Sprint 3-6)"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent_purple']
        p.space_after = Pt(12)
        
        shortterm_items = [
            "Multi-chart dashboards",
            "Saved visualization templates",
            "Anomaly detection alerts",
            "Natural language insights generation",
            "Real-time collaboration features",
        ]
        
        for item in shortterm_items:
            p = shortterm_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text_dark']
            p.space_after = Pt(8)
    
    def add_closing_slide(self, prs: Presentation):
        """Final slide with call to action using Conexus closing template"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['closing']])
        
        # Update any text placeholders if they exist
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "Thank You"
                p.font.size = Pt(54)
                p.font.bold = True
                p.font.color.rgb = self.colors['azure_blue']
                p.alignment = PP_ALIGN.CENTER
        
        # Add custom thank you text
        text_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5), Inches(9.333), Inches(2.0)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = "Ready to transform your data analysis experience"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(26)
        p.font.color.rgb = self.colors['text_dark']
        p.space_after = Pt(16)
        
        p2 = text_frame.add_paragraph()
        p2.text = "NL2Q Analytics Team • October 2025"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(18)
        p2.font.color.rgb = self.colors['azure_light']
    
    # Helper methods for modern visual styling
    def _add_modern_card(self, slide, left: float, top: float, width: float, height: float,
                         title: str, content: str, color: RGBColor, icon: str = ""):
        """Add modern card with shadow and accent border"""
        # Shadow effect
        shadow = slide.shapes.add_shape(
            1, Inches(left + 0.05), Inches(top + 0.05),
            Inches(width), Inches(height)
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = self.colors['gray_200']
        shadow.fill.transparency = 0.5
        shadow.line.fill.background()
        
        # Main card
        card = slide.shapes.add_shape(
            1, Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors['white']
        card.line.color.rgb = color
        card.line.width = Pt(2.5)
        
        # Accent bar at top
        accent = slide.shapes.add_shape(
            1, Inches(left), Inches(top),
            Inches(width), Inches(0.1)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.2),
            Inches(width - 0.4), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = f"{icon} {title}" if icon else title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.65),
            Inches(width - 0.4), Inches(height - 0.85)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = content
        p.font.size = Pt(13)
        p.font.color.rgb = self.colors['text_dark']
        
    def _add_bullet_point(self, slide, icon: str, text: str, left: Inches, top: Inches):
        """Add bullet point with icon"""
        text_box = slide.shapes.add_textbox(left, top, Inches(9.5), Inches(0.7))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = f"{icon}  {text}"
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors['text_dark']
        
    def _add_feature_box(self, slide, title: str, desc: str, left: Inches, top: Inches, width: Inches, height: Inches):
        """Add styled feature box"""
        # Background shape
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['white']
        shape.line.color.rgb = self.colors['primary_light']
        shape.line.width = Pt(2)
        
        # Text
        text_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.2),
            width - Inches(0.4), height - Inches(0.4)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.space_after = Pt(6)
        
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = self.colors['text_dark']
        
    def _add_architecture_layer(self, slide, layer: dict, left: Inches, top: Inches, width: Inches, height: Inches):
        """Add architecture layer box"""
        # Background
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['gray_100']
        shape.line.color.rgb = self.colors['primary']
        shape.line.width = Pt(1.5)
        
        # Title
        title_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.1),
            width - Inches(0.4), Inches(0.3)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = layer['title']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['azure_dark']
        
        # Tech stack
        tech_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.4),
            width - Inches(0.4), Inches(0.25)
        )
        tech_frame = tech_box.text_frame
        p = tech_frame.paragraphs[0]
        p.text = layer['tech']
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['accent_purple']
        
        # Components
        comp_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.7),
            width - Inches(0.4), Inches(0.3)
        )
        comp_frame = comp_box.text_frame
        p = comp_frame.paragraphs[0]
        p.text = " • ".join(layer['components'])
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors['text_dark']
        
    def _add_feature_card(self, slide, feature: dict, left: Inches, top: Inches, width: Inches, height: Inches):
        """Add feature card with status"""
        # Status color mapping
        status_colors = {
            "Delivered": self.colors['success'],
            "In Progress": self.colors['primary'],
            "Planned": self.colors['text_light']
        }
        
        # Background
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['white']
        shape.line.color.rgb = status_colors.get(feature['status'], self.colors['text_medium'])
        shape.line.width = Pt(2)
        
        # Content
        text_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.15),
            width - Inches(0.3), height - Inches(0.3)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = f"{feature['icon']} {feature['title']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['text_dark']
        p.space_after = Pt(4)
        
        p2 = text_frame.add_paragraph()
        p2.text = feature['desc']
        p2.font.size = Pt(11)
        p2.font.color.rgb = self.colors['text_dark']
        p2.space_after = Pt(4)
        
        p3 = text_frame.add_paragraph()
        p3.text = f"Status: {feature['status']}"
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = status_colors.get(feature['status'], self.colors['text_medium'])
        
    def _add_flow_step(self, slide, icon: str, title: str, desc: str, left: Inches, top: Inches, width: Inches):
        """Add flow step"""
        # Background
        shape = slide.shapes.add_shape(1, left, top, width, Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['white']
        shape.line.color.rgb = self.colors['primary_light']
        shape.line.width = Pt(1.5)
        
        # Content
        text_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.15),
            width - Inches(0.4), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = f"{icon} {title}: {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['text_dark']
        
    def _add_metric_card(self, slide, value: str, title: str, desc: str, left: Inches, top: Inches, width: Inches, height: Inches):
        """Add metric card"""
        # Background with gradient effect
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # Value
        value_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.3),
            width - Inches(0.2), Inches(0.6)
        )
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = value
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # Title
        title_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.95),
            width - Inches(0.2), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # Description
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(1.35),
            width - Inches(0.2), Inches(0.5)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = desc
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors['primary_light']
        
    def _add_modern_metric(self, slide, left: float, top: float, width: float, height: float,
                           value: str, label: str, sublabel: str, color: RGBColor):
        """Add modern metric card with gradient accent"""
        # Shadow
        shadow = slide.shapes.add_shape(
            1, Inches(left + 0.03), Inches(top + 0.03),
            Inches(width), Inches(height)
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = self.colors['gray_200']
        shadow.fill.transparency = 0.6
        shadow.line.fill.background()
        
        # Main card with gradient
        card = slide.shapes.add_shape(
            1, Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        # Value text
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.3),
            Inches(width - 0.2), Inches(0.6)
        )
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = value
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.95),
            Inches(width - 0.2), Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.word_wrap = True
        p = label_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # Sublabel
        sublabel_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 1.35),
            Inches(width - 0.2), Inches(0.4)
        )
        sublabel_frame = sublabel_box.text_frame
        sublabel_frame.word_wrap = True
        p = sublabel_frame.paragraphs[0]
        p.text = sublabel
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['gray_100']
    
    def generate_mvp_deck(self) -> str:
        """Generate complete MVP deck with Conexus template"""
        print("🎨 Generating NL2Q Analyst MVP Deck with Conexus Corporate Template...")
        
        prs = self.create_presentation()
        
        # Slide 1: Title
        print("  📄 Adding title slide...")
        self.add_title_slide(prs)
        
        # Section Divider: Overview
        print("  📄 Adding section divider: Overview...")
        self.add_section_divider(prs, "Overview")
        
        # Slide 2: Problem Statement
        print("  📄 Adding problem statement...")
        self.add_problem_statement_slide(prs)
        
        # Slide 3: Solution
        print("  📄 Adding solution overview...")
        self.add_solution_slide(prs)
        
        # Slide 4: Architecture
        print("  📄 Adding technical architecture...")
        self.add_architecture_slide(prs)
        
        # Section Divider: Demo
        print("  📄 Adding section divider: Demo...")
        self.add_section_divider(prs, "Live Demo", use_orange=True)
        
        # Slides 5-8: Screenshots with context
        screenshot_contexts = [
            ("Feature Demo: Comparison Query", "Natural language query automatically generates comparison-based KPIs and contextual insights"),
            ("Feature Demo: Adaptive Layout", "Intelligent layout with filtered KPIs, dynamic charts, and interactive comparison cards"),
            ("Feature Demo: Contextual Insights", "Priority-ranked insights sidebar with click-to-filter functionality"),
            ("Feature Demo: Professional UI", "Enterprise-grade design with optimized spacing and color scheme"),
        ]
        
        for idx, (title, desc) in enumerate(screenshot_contexts):
            screenshot_file = self.assets_dir / self.screenshots[idx]
            print(f"  📄 Adding screenshot slide {idx+1}: {screenshot_file.name}...")
            self.add_screenshot_slide(prs, screenshot_file, title, desc)
        
        # Section Divider: Features & Results
        print("  📄 Adding section divider: Features...")
        self.add_section_divider(prs, "Features & Results")
        
        # Slide 9: Key Features
        print("  📄 Adding key features...")
        self.add_key_features_slide(prs)
        
        # Slide 10: Demo Flow
        print("  📄 Adding demo flow...")
        self.add_demo_flow_slide(prs)
        
        # Slide 11: Results
        print("  📄 Adding results & impact...")
        self.add_results_slide(prs)
        
        # Section Divider: Roadmap
        print("  📄 Adding section divider: Roadmap...")
        self.add_section_divider(prs, "Next Steps", use_orange=True)
        
        # Slide 12: Next Steps
        print("  📄 Adding next steps...")
        self.add_next_steps_slide(prs)
        
        # Slide 13: Closing
        print("  📄 Adding closing slide...")
        self.add_closing_slide(prs)
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"nl2q_analyst_mvp_conexus_{timestamp}.pptx"
        filepath = self.output_dir / filename
        
        prs.save(str(filepath))
        print(f"✅ MVP deck generated: {filepath}")
        print(f"📊 Total slides: {len(prs.slides)}")
        print(f"🎨 Using Conexus Corporate Template with modern design")
        
        return str(filepath)


def main():
    """Generate NL2Q Analyst MVP presentation"""
    try:
        generator = NL2QMVPDeckGenerator(
            assets_dir="assets",
            output_dir="outputs"
        )
        
        filepath = generator.generate_mvp_deck()
        
        print("\n" + "="*70)
        print("🎉 NL2Q Analyst MVP Deck Complete!")
        print("="*70)
        print(f"📁 Location: {filepath}")
        print(f"🎨 Template: Conexus Corporate 2025 with Modern Design")
        print(f"📊 Total Slides: {len(filepath.split('/')[-1])}")
        print("\n📋 Deck Structure:")
        print("   • Title Slide (Conexus branded)")
        print("   • Section: Overview")
        print("     - Problem Statement")
        print("     - Solution Overview")
        print("     - Technical Architecture")
        print("   • Section: Live Demo")
        print("     - 4 Feature Screenshots with context")
        print("   • Section: Features & Results")
        print("     - Key Features Delivered")
        print("     - Demo Flow")
        print("     - Results & Impact")
        print("   • Section: Next Steps")
        print("     - Roadmap & Next Steps")
        print("   • Closing Slide (Conexus branded)")
        print("\n✨ Features:")
        print("   ✓ Professional Conexus corporate template")
        print("   ✓ Modern visual design with cards and shadows")
        print("   ✓ Section dividers for clear organization")
        print("   ✓ Consistent branding throughout")
        print("\n💡 Ready for executive presentations!")
        
    except Exception as e:
        print(f"❌ Error generating deck: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
