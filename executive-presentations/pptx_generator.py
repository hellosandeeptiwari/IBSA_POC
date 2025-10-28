"""
Enterprise-Grade PowerPoint (PPTX) Presentation Generator
Converts content to professional PowerPoint decks

Author: NL2Q Analytics Team
Date: October 2025
Version: 2.0 - Enhanced Visual Design
"""

import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from typing import Dict, List, Optional, Tuple


class EnterprisePPTXGenerator:
    """
    Generate high-end consulting-style PowerPoint presentations
    Consistent with McKinsey, BCG, Bain presentation standards
    """
    
    def __init__(self, output_dir: str = "outputs"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Enterprise color palette (RGB values)
        self.colors = {
            "primary": RGBColor(26, 26, 46),       # Deep navy
            "secondary": RGBColor(22, 33, 62),      # Dark blue
            "accent": RGBColor(15, 52, 96),         # Royal blue
            "highlight": RGBColor(83, 52, 131),     # Purple
            "success": RGBColor(45, 106, 79),       # Forest green
            "warning": RGBColor(244, 162, 97),      # Warm orange
            "danger": RGBColor(214, 40, 40),        # Professional red
            "light": RGBColor(248, 249, 250),       # Off-white
            "muted": RGBColor(108, 117, 125),       # Gray
            "white": RGBColor(255, 255, 255),
            "text": RGBColor(33, 37, 41)
        }
        
    def create_presentation(self) -> Presentation:
        """Create a blank presentation with standard dimensions"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        return prs
        
    def add_title_slide(self, prs: Presentation, title: str, subtitle: str, 
                       company: str = "NL2Q Analytics") -> None:
        """Add professional title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background gradient effect (solid color approximation)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Inter'
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2), Inches(11.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.font.name = 'Inter'
        
        # Company/Date footer
        footer_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(11.333), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = f"{company} | {datetime.now().strftime('%B %Y')}"
        
        p = footer_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.font.name = 'Inter'
        
    def add_section_slide(self, prs: Presentation, section_title: str, 
                         section_number: str = "") -> None:
        """Add section divider slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['accent']
        
        # Section number (top)
        if section_number:
            num_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(11.333), Inches(0.8)
            )
            num_frame = num_box.text_frame
            num_frame.text = section_number
            
            p = num_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(150, 150, 200)
            p.font.name = 'Inter'
        
        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.2), Inches(11.333), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Inter'
        
    def add_content_slide(self, prs: Presentation, title: str, 
                         content: List[Dict[str, any]], slide_number: int = 0) -> None:
        """Add content slide with title and bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header bar
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(13.333), Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['primary']
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(10), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Inter'
        
        # Slide number
        if slide_number > 0:
            num_box = slide.shapes.add_textbox(
                Inches(11.8), Inches(0.15), Inches(1), Inches(0.5)
            )
            num_frame = num_box.text_frame
            num_frame.text = str(slide_number)
            
            p = num_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(180, 180, 180)
            p.font.name = 'Inter'
        
        # Content area
        content_top = 1.2
        content_height = 6.0
        current_y = content_top
        
        for item in content:
            if item['type'] == 'bullets':
                self._add_bullet_list(slide, item['items'], 
                                     Inches(1.2), Inches(current_y),
                                     Inches(11), Inches(content_height))
            elif item['type'] == 'two_column':
                self._add_two_column(slide, item['left'], item['right'],
                                    Inches(current_y))
            elif item['type'] == 'metrics':
                self._add_metrics(slide, item['metrics'], current_y)
                current_y += 2.5  # Space after metrics
            elif item['type'] == 'table':
                self._add_table(slide, item['headers'], item['rows'],
                               Inches(current_y))
                               
    def _add_bullet_list(self, slide, items: List[str], left, top, width, height):
        """Add formatted bullet point list"""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(items):
            if i > 0:
                text_frame.add_paragraph()
            
            p = text_frame.paragraphs[i]
            
            # Parse bullet level (indentation)
            level = 0
            text = item
            if item.startswith('  '):
                level = 1
                text = item.strip()
            
            p.text = text
            p.level = level
            p.font.size = Pt(18 if level == 0 else 16)
            p.font.color.rgb = self.colors['text']
            p.font.name = 'Inter'
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            
    def _add_two_column(self, slide, left_items: List[str], 
                       right_items: List[str], top):
        """Add two-column layout"""
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(1.2), top, Inches(5.2), Inches(5.5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, item in enumerate(left_items):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text']
            p.font.name = 'Inter'
            p.space_after = Pt(10)
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(6.9), top, Inches(5.2), Inches(5.5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, item in enumerate(right_items):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text']
            p.font.name = 'Inter'
            p.space_after = Pt(10)
            
    def _add_metrics(self, slide, metrics: List[Dict], top):
        """Add metric cards"""
        num_metrics = len(metrics)
        card_width = 2.4
        spacing = 0.3
        total_width = (card_width * num_metrics) + (spacing * (num_metrics - 1))
        start_left = (13.333 - total_width) / 2
        
        # Convert top to float if it's Inches object
        top_value = top / Inches(1) if hasattr(top, '__truediv__') else top
        
        for i, metric in enumerate(metrics):
            left = start_left + (i * (card_width + spacing))
            
            # Card background
            card = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(left), Inches(top_value), Inches(card_width), Inches(1.8)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors['light']
            card.line.color.rgb = self.colors['accent']
            card.line.width = Pt(2)
            
            # Metric value
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top_value + 0.3), 
                Inches(card_width - 0.4), Inches(0.7)
            )
            value_frame = value_box.text_frame
            value_frame.text = metric['value']
            
            p = value_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.colors['accent']
            p.font.name = 'Inter'
            
            # Metric label
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top_value + 1.1),
                Inches(card_width - 0.4), Inches(0.5)
            )
            label_frame = label_box.text_frame
            label_frame.text = metric['label']
            label_frame.word_wrap = True
            
            p = label_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['muted']
            p.font.name = 'Inter'
            
    def _add_table(self, slide, headers: List[str], rows: List[List[str]], top):
        """Add formatted table"""
        # Calculate dimensions
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header
        table_width = 11
        table_height = min(5, num_rows * 0.5)
        
        # Add table shape
        left = Inches(1.2)
        width = Inches(table_width)
        height = Inches(table_height)
        
        table = slide.shapes.add_table(
            num_rows, num_cols, left, top, width, height
        ).table
        
        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
            
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.font.name = 'Inter'
            p.alignment = PP_ALIGN.CENTER
        
        # Data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                
                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(250, 250, 250)
                
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = self.colors['text']
                p.font.name = 'Inter'
                p.alignment = PP_ALIGN.CENTER
                
    def generate_powerbi_integration_deck(self) -> str:
        """Generate complete Power BI Integration & Multi-Tenancy deck"""
        print("🎨 Generating Enterprise PowerPoint Presentation...")
        
        prs = self.create_presentation()
        
        # Slide 1: Title
        self.add_title_slide(
            prs,
            "Power BI Integration Strategy",
            "NL2Q Analyst Platform | Multi-Tenant Architecture"
        )
        
        # Slide 2: Executive Summary
        self.add_content_slide(prs, "Executive Summary", [
            {
                'type': 'bullets',
                'items': [
                    "Strategic Integration: Power BI will consume NL2Q Analyst's natural language query capabilities",
                    "Three Implementation Approaches: Custom Connector, REST API Gateway, and Embedded Analytics",
                    "Multi-Tenant Architecture: Schema-per-tenant isolation with comprehensive security",
                    "Timeline: 12-16 weeks for full implementation",
                    "Investment: $180K-$280K depending on approach",
                    "Business Impact: 40% reduction in report creation time, enhanced data accessibility"
                ]
            }
        ], 2)
        
        # Slide 3: Section - Architecture
        self.add_section_slide(prs, "Technical Architecture", "01")
        
        # Slide 4: Integration Approaches
        self.add_content_slide(prs, "Three Integration Approaches", [
            {
                'type': 'two_column',
                'left': [
                    "🔌 Approach 1: Custom Connector",
                    "• Power Query M language",
                    "• Native Power BI experience",
                    "• Full data modeling support",
                    "• Timeline: 12-16 weeks",
                    "• Cost: $180K-$220K"
                ],
                'right': [
                    "🌐 Approach 2: REST API Gateway",
                    "• OData/JSON endpoints",
                    "• Flexible integration",
                    "• Web connector compatible",
                    "• Timeline: 8-10 weeks",
                    "• Cost: $120K-$160K"
                ]
            }
        ], 4)
        
        # Slide 5: Multi-Tenant Architecture
        self.add_content_slide(prs, "Multi-Tenant Security Model", [
            {
                'type': 'bullets',
                'items': [
                    "🔐 Three-Tier Isolation Strategy",
                    "  Tier 1: Schema-per-tenant in Snowflake (physical separation)",
                    "  Tier 2: Row-Level Security policies (logical separation)",
                    "  Tier 3: JWT authentication with tenant context",
                    "🛡️ Security Features",
                    "  Azure Key Vault for credential management",
                    "  Role-based access control (RBAC)",
                    "  Audit logging and compliance tracking",
                    "📊 Performance Optimization",
                    "  Redis caching per tenant",
                    "  Connection pooling",
                    "  Query result caching"
                ]
            }
        ], 5)
        
        # Slide 6: Implementation Timeline
        self.add_content_slide(prs, "Implementation Roadmap", [
            {
                'type': 'table',
                'headers': ['Phase', 'Duration', 'Key Deliverables', 'Investment'],
                'rows': [
                    ['Phase 1: Foundation', '4 weeks', 'Multi-tenant auth, schema isolation', '$45K-$55K'],
                    ['Phase 2: API Development', '3 weeks', 'OData endpoints, Power BI connector', '$35K-$45K'],
                    ['Phase 3: Integration', '3 weeks', 'Custom connector, testing', '$40K-$50K'],
                    ['Phase 4: Deployment', '2 weeks', 'Production rollout, monitoring', '$30K-$40K'],
                    ['Contingency', '2 weeks', 'Buffer for issues', '$30K-$40K']
                ]
            }
        ], 6)
        
        # Slide 7: Section - Cost Analysis
        self.add_section_slide(prs, "Investment & Returns", "02")
        
        # Slide 8: Cost Breakdown
        self.add_content_slide(prs, "Investment Breakdown", [
            {
                'type': 'metrics',
                'metrics': [
                    {'value': '$180K', 'label': 'Total Implementation'},
                    {'value': '$40K/yr', 'label': 'Infrastructure Costs'},
                    {'value': '12 weeks', 'label': 'Time to Market'},
                    {'value': '40%', 'label': 'Efficiency Gain'}
                ]
            }
        ], 8)
        
        # Slide 9: Technical Feasibility
        self.add_content_slide(prs, "Technical Feasibility Assessment", [
            {
                'type': 'two_column',
                'left': [
                    "✅ Strengths",
                    "• Existing FastAPI backend",
                    "• GPT-4o/o3-mini proven",
                    "• Snowflake production-ready",
                    "• React frontend mature",
                    "• Pinecone semantic search"
                ],
                'right': [
                    "⚠️ Gaps to Address",
                    "• No authentication system",
                    "• No multi-tenant isolation",
                    "• Missing audit logging",
                    "• No monitoring/alerting",
                    "• Power BI connector needed"
                ]
            }
        ], 9)
        
        # Slide 10: Risk Assessment
        self.add_content_slide(prs, "Risk Analysis & Mitigation", [
            {
                'type': 'table',
                'headers': ['Risk', 'Impact', 'Probability', 'Mitigation'],
                'rows': [
                    ['Data Isolation Breach', 'Critical', 'Low', 'Schema-per-tenant + RLS + audit'],
                    ['Performance Degradation', 'High', 'Medium', 'Redis caching + connection pooling'],
                    ['API Rate Limits', 'Medium', 'Medium', 'Token bucket + tenant quotas'],
                    ['Power BI Connector Issues', 'High', 'Low', 'Thorough testing + fallback to REST'],
                    ['Cost Overruns', 'Medium', 'Medium', '20% contingency buffer']
                ]
            }
        ], 10)
        
        # Slide 11: Section - Benefits
        self.add_section_slide(prs, "Business Impact", "03")
        
        # Slide 12: Business Benefits
        self.add_content_slide(prs, "Expected Business Outcomes", [
            {
                'type': 'bullets',
                'items': [
                    "📈 Productivity Gains",
                    "  40% reduction in report creation time",
                    "  60% faster ad-hoc query responses",
                    "  70% decrease in SQL developer dependency",
                    "🎯 User Experience",
                    "  Natural language queries in familiar Power BI interface",
                    "  Self-service analytics for business users",
                    "  Reduced training requirements",
                    "💰 Financial Impact",
                    "  $250K annual savings in analyst time",
                    "  $180K implementation cost",
                    "  18-month ROI break-even"
                ]
            }
        ], 12)
        
        # Slide 13: Competitive Advantage
        self.add_content_slide(prs, "Market Differentiation", [
            {
                'type': 'two_column',
                'left': [
                    "🚀 Unique Capabilities",
                    "• AI-powered SQL generation",
                    "• Semantic schema search",
                    "• Multi-database support",
                    "• Enterprise-grade security",
                    "• Native Power BI integration"
                ],
                'right': [
                    "🏆 vs. Competitors",
                    "• Microsoft: No NL2SQL in Power BI",
                    "• Tableau: Limited AI capabilities",
                    "• Looker: No custom connector",
                    "• ThoughtSpot: 3x higher cost",
                    "• Seekwell: No multi-tenancy"
                ]
            }
        ], 13)
        
        # Slide 14: Recommendations
        self.add_content_slide(prs, "Strategic Recommendations", [
            {
                'type': 'bullets',
                'items': [
                    "🎯 Phase 1 Priority: Implement REST API Gateway approach (fastest time-to-market)",
                    "📅 Timeline: Begin Q1 2025, production-ready by Q2 2025",
                    "💼 Resource Allocation: 3 full-stack engineers, 1 DevOps engineer",
                    "🔐 Security First: Complete authentication & multi-tenancy before Power BI integration",
                    "📊 Pilot Program: 3 enterprise customers for beta testing (8 weeks)",
                    "📈 Success Metrics: Track query accuracy, response time, user adoption"
                ]
            }
        ], 14)
        
        # Slide 15: Closing
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']
        
        # Main text
        closing_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.8), Inches(9.333), Inches(2)
        )
        closing_frame = closing_box.text_frame
        closing_frame.text = "Questions & Discussion"
        
        p = closing_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Inter'
        
        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(2), Inches(5.2), Inches(9.333), Inches(1)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "NL2Q Analytics Team\nenterprise@nl2q-analytics.com"
        
        for p in contact_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(180, 180, 180)
            p.font.name = 'Inter'
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"nl2q_powerbi_strategy_{timestamp}.pptx"
        filepath = self.output_dir / filename
        
        prs.save(str(filepath))
        
        print(f"\n✅ PowerPoint Presentation generated successfully!")
        print(f"📄 File: {filepath}")
        print(f"📊 Slides: {len(prs.slides)}")
        print(f"🎯 Topics: Power BI Integration, Multi-Tenancy, Timeline, Costs, Risks, Benefits")
        
        return str(filepath)


def main():
    """Main execution"""
    generator = EnterprisePPTXGenerator()
    output_file = generator.generate_powerbi_integration_deck()
    
    print(f"\n🎉 Presentation ready: {output_file}")
    print("\n💡 Next steps:")
    print("1. Open the PPTX file in PowerPoint")
    print("2. Review and customize content")
    print("3. Adjust company branding/colors")
    print("4. Export as PDF if needed")


if __name__ == "__main__":
    main()
