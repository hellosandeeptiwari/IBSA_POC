"""Extract text content from reference presentation"""
from pptx import Presentation
from pathlib import Path

ref_path = Path("assets/Quarterly Business Summary 7.2025 v3-forConexus 1.pptx")
prs = Presentation(ref_path)

print(f"\n{'='*80}")
print(f"REFERENCE PRESENTATION CONTENT EXTRACTION")
print(f"{'='*80}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*80}")
    print(f"SLIDE {i}")
    print(f"{'='*80}")
    
    # Extract all text from slide
    all_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            all_text.append(shape.text.strip())
    
    if all_text:
        for text in all_text:
            print(text)
            print("-" * 40)
    else:
        print("(No text content)")
    
    # Check for tables
    for shape in slide.shapes:
        if shape.shape_type == 19:  # Table
            print("\n📋 TABLE CONTENT:")
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                print(f"  Row {row_idx}: {row_text}")
