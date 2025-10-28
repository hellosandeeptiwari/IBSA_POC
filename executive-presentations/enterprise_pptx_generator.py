"""
Enterprise-Grade PowerPoint Generator - Premium Edition
Professional decks with stunning visuals and sophisticated design

Author: NL2Q Analytics Team
Date: October 2025
Version: 2.0 - Premium Visual Experience
"""

import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from typing import Dict, List, Optional, Tuple


class EnterprisePremiumPPTX:
    """
    Generate stunning enterprise presentations with premium visual design
    """
    
    def __init__(self, output_dir: str = "outputs"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Premium color palette (sophisticated, modern)
        self.colors = {
            # Primary gradient colors
            "primary_dark": RGBColor(15, 23, 42),      # Slate 900
            "primary": RGBColor(30, 41, 59),           # Slate 800
            "primary_light": RGBColor(51, 65, 85),     # Slate 700
            
            # Accent colors (vibrant, professional)
            "accent_blue": RGBColor(59, 130, 246),     # Blue 500
            "accent_purple": RGBColor(139, 92, 246),   # Purple 500
            "accent_cyan": RGBColor(6, 182, 212),      # Cyan 500
            "accent_emerald": RGBColor(16, 185, 129),  # Emerald 500
            
            # Status colors
            "success": RGBColor(34, 197, 94),          # Green 500
            "warning": RGBColor(251, 146, 60),         # Orange 400
            "info": RGBColor(96, 165, 250),            # Blue 400
            
            # Neutral colors
            "white": RGBColor(255, 255, 255),
            "gray_50": RGBColor(249, 250, 251),
            "gray_100": RGBColor(243, 244, 246),
            "gray_200": RGBColor(229, 231, 235),
            "gray_300": RGBColor(209, 213, 219),
            "gray_600": RGBColor(75, 85, 99),
            "gray_700": RGBColor(55, 65, 81),
            "gray_900": RGBColor(17, 24, 39),
            
            # Text colors
            "text_dark": RGBColor(30, 41, 59),
            "text_medium": RGBColor(71, 85, 105),
            "text_light": RGBColor(148, 163, 184),
        }
        
    def create_presentation(self) -> Presentation:
        """Create presentation with widescreen dimensions"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 widescreen
        prs.slide_height = Inches(7.5)
        return prs
    
    def add_gradient_background(self, slide, color1: RGBColor, color2: RGBColor):
        """Add sophisticated gradient background"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45.0
        
        # First gradient stop
        fill.gradient_stops[0].color.rgb = color1
        
        # Second gradient stop
        fill.gradient_stops[1].color.rgb = color2
        
    def add_decorative_elements(self, slide):
        """Add modern decorative shapes for visual interest"""
        # Top-right accent circle
        circle = slide.shapes.add_shape(
            3,  # Oval
            Inches(11.5), Inches(-0.5), Inches(3), Inches(3)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.colors['accent_cyan']
        circle.fill.transparency = 0.85
        circle.line.fill.background()
        
        # Bottom-left accent circle
        circle2 = slide.shapes.add_shape(
            3,  # Oval
            Inches(-1), Inches(5.5), Inches(2.5), Inches(2.5)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = self.colors['accent_purple']
        circle2.fill.transparency = 0.85
        circle2.line.fill.background()
        
    def add_premium_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """Create stunning title slide with modern design"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Gradient background
        self.add_gradient_background(slide, self.colors['primary_dark'], self.colors['primary'])
        
        # Decorative elements
        self.add_decorative_elements(slide)
        
        # Large accent bar (left side) - Azure blue
        accent_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(0.3), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors['accent_cyan']
        accent_bar.line.fill.background()
        
        # Main title with modern styling
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.2), Inches(10), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(64)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Segoe UI'
        p.line_spacing = 1.1
        
        # Subtitle with accent color
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4.5), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(28)
        p.font.color.rgb = self.colors['accent_cyan']
        p.font.name = 'Segoe UI Light'
        
        # Azure Partner badge and footer
        date_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(6.5), Inches(10), Inches(0.6)
        )
        date_frame = date_box.text_frame
        date_frame.text = "Microsoft Azure Partner • Built on Azure Foundry"
        date_frame.add_paragraph()
        date_frame.paragraphs[1].text = f"NL2Q Analytics • {datetime.now().strftime('%B %Y')}"
        
        for p in date_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.color.rgb = self.colors['gray_300']
            p.font.name = 'Segoe UI'
        
    def add_section_divider(self, prs: Presentation, section_title: str, 
                           section_subtitle: str = "", icon_text: str = ""):
        """Add modern section divider with visual impact"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Gradient background
        self.add_gradient_background(slide, self.colors['primary'], self.colors['primary_light'])
        
        # Large accent shape (diagonal)
        accent = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(8), Inches(-1), Inches(7), Inches(9)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.colors['accent_purple']
        accent.fill.transparency = 0.9
        accent.line.fill.background()
        accent.rotation = 15.0
        
        # Section icon/number
        if icon_text:
            icon_box = slide.shapes.add_textbox(
                Inches(2), Inches(1.5), Inches(2), Inches(1.5)
            )
            icon_frame = icon_box.text_frame
            icon_frame.text = icon_text
            
            p = icon_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = self.colors['accent_cyan']
            p.font.name = 'Segoe UI'
        
        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.2), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Segoe UI'
        
        # Section subtitle
        if section_subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2), Inches(5), Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = section_subtitle
            subtitle_frame.word_wrap = True
            
            p = subtitle_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(22)
            p.font.color.rgb = self.colors['gray_200']
            p.font.name = 'Segoe UI Light'
    
    def add_content_slide(self, prs: Presentation, title: str, 
                         slide_number: int = 0):
        """Create content slide with premium styling"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background with subtle texture
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['white']
        
        # Decorative top accent
        top_accent = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(13.333), Inches(0.05)
        )
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = self.colors['accent_cyan']
        top_accent.line.fill.background()
        
        # Header area with gradient
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0.05), Inches(13.333), Inches(1.1)
        )
        header.fill.gradient()
        header.fill.gradient_angle = 0
        header.fill.gradient_stops[0].color.rgb = self.colors['primary_dark']
        header.fill.gradient_stops[1].color.rgb = self.colors['primary']
        header.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.25), Inches(11), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Segoe UI'
        
        # Slide number with modern styling
        if slide_number > 0:
            num_circle = slide.shapes.add_shape(
                3,  # Oval
                Inches(12.1), Inches(0.3), Inches(0.6), Inches(0.6)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = self.colors['accent_cyan']
            num_circle.line.fill.background()
            
            num_box = slide.shapes.add_textbox(
                Inches(12.1), Inches(0.3), Inches(0.6), Inches(0.6)
            )
            num_frame = num_box.text_frame
            num_frame.text = str(slide_number)
            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = num_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.font.name = 'Segoe UI'
        
        # Bottom decorative line
        bottom_line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.6), Inches(7.3), Inches(12.133), Inches(0.03)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = self.colors['gray_200']
        bottom_line.line.fill.background()
        
        return slide
    
    def add_icon_bullets(self, slide, items: List[str], left: float, top: float, 
                        icon_color: RGBColor = None):
        """Add bullet points with colorful icons"""
        if icon_color is None:
            icon_color = self.colors['accent_blue']
        
        y_offset = top
        for item in items:
            # Determine indent level
            level = 0
            text = item
            if item.startswith('  '):
                level = 1
                text = item.strip()
                icon_color = self.colors['accent_purple']
            
            indent = level * 0.5
            
            # Icon/bullet (modern square)
            bullet = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(left + indent), Inches(y_offset), Inches(0.15), Inches(0.15)
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = icon_color
            bullet.line.fill.background()
            
            # Text
            text_box = slide.shapes.add_textbox(
                Inches(left + indent + 0.3), Inches(y_offset - 0.05),
                Inches(10.5 - indent), Inches(0.5)
            )
            text_frame = text_box.text_frame
            text_frame.text = text
            text_frame.word_wrap = True
            
            p = text_frame.paragraphs[0]
            p.font.size = Pt(20 if level == 0 else 18)
            p.font.color.rgb = self.colors['text_dark'] if level == 0 else self.colors['text_medium']
            p.font.name = 'Segoe UI'
            p.line_spacing = 1.3
            
            y_offset += 0.45 if level == 0 else 0.4
    
    def add_premium_metrics(self, slide, metrics: List[Dict], top: float):
        """Add stunning metric cards with modern design"""
        num_metrics = len(metrics)
        card_width = 2.6
        spacing = 0.4
        total_width = (card_width * num_metrics) + (spacing * (num_metrics - 1))
        start_left = (13.333 - total_width) / 2
        
        colors = [
            self.colors['accent_blue'],
            self.colors['accent_purple'],
            self.colors['accent_cyan'],
            self.colors['accent_emerald']
        ]
        
        for i, metric in enumerate(metrics):
            left = start_left + (i * (card_width + spacing))
            color = colors[i % len(colors)]
            
            # Card shadow (subtle depth)
            shadow = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(left + 0.05), Inches(top + 0.05), 
                Inches(card_width), Inches(2.2)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = self.colors['gray_200']
            shadow.fill.transparency = 0.5
            shadow.line.fill.background()
            
            # Main card
            card = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(left), Inches(top), Inches(card_width), Inches(2.2)
            )
            card.fill.gradient()
            card.fill.gradient_angle = 135
            card.fill.gradient_stops[0].color.rgb = self.colors['white']
            card.fill.gradient_stops[1].color.rgb = self.colors['gray_50']
            card.line.color.rgb = color
            card.line.width = Pt(3)
            
            # Accent top bar
            accent_bar = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(left), Inches(top), Inches(card_width), Inches(0.15)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = color
            accent_bar.line.fill.background()
            
            # Metric value with emphasis
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 0.5),
                Inches(card_width - 0.4), Inches(0.9)
            )
            value_frame = value_box.text_frame
            value_frame.text = metric['value']
            value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = value_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = color
            p.font.name = 'Segoe UI'
            
            # Metric label
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 1.5),
                Inches(card_width - 0.4), Inches(0.6)
            )
            label_frame = label_box.text_frame
            label_frame.text = metric['label']
            label_frame.word_wrap = True
            
            p = label_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text_medium']
            p.font.name = 'Segoe UI'
            p.line_spacing = 1.2
    
    def add_two_column_premium(self, slide, left_title: str, left_items: List[str],
                               right_title: str, right_items: List[str], top: float):
        """Add premium two-column layout with modern cards"""
        # Left column card
        left_card = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(top), Inches(5.6), Inches(5.2)
        )
        left_card.fill.gradient()
        left_card.fill.gradient_angle = 45
        left_card.fill.gradient_stops[0].color.rgb = self.colors['gray_50']
        left_card.fill.gradient_stops[1].color.rgb = self.colors['white']
        left_card.line.color.rgb = self.colors['accent_blue']
        left_card.line.width = Pt(2)
        
        # Left title bar
        left_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(top), Inches(5.6), Inches(0.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = self.colors['accent_blue']
        left_bar.line.fill.background()
        
        # Left title
        left_title_box = slide.shapes.add_textbox(
            Inches(1), Inches(top + 0.05), Inches(5.2), Inches(0.4)
        )
        left_title_frame = left_title_box.text_frame
        left_title_frame.text = left_title
        left_title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = left_title_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Segoe UI'
        
        # Left content
        self.add_icon_bullets(slide, left_items, 1.1, top + 0.7, self.colors['accent_blue'])
        
        # Right column card
        right_card = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(6.9), Inches(top), Inches(5.6), Inches(5.2)
        )
        right_card.fill.gradient()
        right_card.fill.gradient_angle = 45
        right_card.fill.gradient_stops[0].color.rgb = self.colors['gray_50']
        right_card.fill.gradient_stops[1].color.rgb = self.colors['white']
        right_card.line.color.rgb = self.colors['accent_purple']
        right_card.line.width = Pt(2)
        
        # Right title bar
        right_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(6.9), Inches(top), Inches(5.6), Inches(0.5)
        )
        right_bar.fill.solid()
        right_bar.fill.fore_color.rgb = self.colors['accent_purple']
        right_bar.line.fill.background()
        
        # Right title
        right_title_box = slide.shapes.add_textbox(
            Inches(7.1), Inches(top + 0.05), Inches(5.2), Inches(0.4)
        )
        right_title_frame = right_title_box.text_frame
        right_title_frame.text = right_title
        right_title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = right_title_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.font.name = 'Segoe UI'
        
        # Right content
        self.add_icon_bullets(slide, right_items, 7.2, top + 0.7, self.colors['accent_purple'])
    
    def add_premium_table(self, slide, headers: List[str], rows: List[List[str]], top: float):
        """Add modern data table with premium styling"""
        num_cols = len(headers)
        num_rows = len(rows) + 1
        
        left = Inches(1)
        width = Inches(11.333)
        height = Inches(min(5, num_rows * 0.6))
        
        table = slide.shapes.add_table(
            num_rows, num_cols, left, Inches(top), width, height
        ).table
        
        # Header row with gradient
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
            
            # Add border
            cell.text_frame.margin_left = Inches(0.1)
            cell.text_frame.margin_right = Inches(0.1)
            
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.CENTER
            
        # Data rows with alternating colors
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                
                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors['gray_50']
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors['white']
                
                cell.text_frame.margin_left = Inches(0.1)
                cell.text_frame.margin_right = Inches(0.1)
                
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.color.rgb = self.colors['text_dark']
                p.font.name = 'Segoe UI'
                p.alignment = PP_ALIGN.CENTER
    
    def add_closing_slide(self, prs: Presentation):
        """Add stunning closing slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Gradient background
        self.add_gradient_background(slide, self.colors['primary_dark'], self.colors['primary'])
        
        # Decorative elements
        self.add_decorative_elements(slide)
        
        # Large text
        closing_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.5), Inches(9.333), Inches(2)
        )
        closing_frame = closing_box.text_frame
        closing_frame.text = "Let's Transform\nData Into Insights"
        closing_frame.word_wrap = True
        
        for para in closing_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(56)
            para.font.bold = True
            para.font.color.rgb = self.colors['white']
            para.font.name = 'Segoe UI'
            para.line_spacing = 1.2
        
        # Accent line
        line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(5), Inches(5), Inches(3.333), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent_cyan']
        line.line.fill.background()
        
        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(2), Inches(5.8), Inches(9.333), Inches(1)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "NL2Q Analytics Team\nReady to discuss your data strategy"
        
        for para in contact_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(22)
            para.font.color.rgb = self.colors['gray_200']
            para.font.name = 'Segoe UI Light'
    
    def generate_powerbi_deck(self) -> str:
        """Generate complete Power BI Integration deck without pricing"""
        print("🎨 Generating Premium Azure-Native Enterprise PowerPoint...")
        
        prs = self.create_presentation()
        
        # Slide 1: Premium Title
        self.add_premium_title_slide(
            prs,
            "Azure-Native Power BI Integration",
            "Built on Azure Foundry | Powered by Azure AI Services"
        )
        
        # Slide 2: Executive Summary
        slide = self.add_content_slide(prs, "Executive Summary", 2)
        self.add_icon_bullets(slide, [
            "☁️ Azure-Native Architecture: Built entirely on Microsoft Azure ecosystem",
            "🤖 Azure AI Services: Powered by Azure OpenAI Service (GPT-4o)",
            "🔧 Azure Foundry Development: Modern cloud-native development platform",
            "🔒 Azure Security: Multi-tenant architecture with Azure AD & Key Vault",
            "⚡ Power BI Integration: Seamless connectivity to Microsoft Power BI",
            "📊 Azure Data Platform: Supports Azure SQL, Synapse, Snowflake on Azure"
        ], 0.8, 1.8)
        
        # Slide 3: Section - Azure Architecture
        self.add_section_divider(
            prs,
            "Azure-Native Architecture",
            "Leveraging Microsoft Azure ecosystem for enterprise-grade analytics",
            "01"
        )
        
        # Slide 4: Azure Technology Stack
        slide = self.add_content_slide(prs, "Azure Technology Stack", 4)
        self.add_two_column_premium(
            slide,
            "☁️ Azure Core Services",
            [
                "Azure OpenAI Service (GPT-4o, GPT-o3-mini)",
                "Azure App Service (FastAPI backend)",
                "Azure Static Web Apps (React frontend)",
                "Azure Container Registry & AKS",
                "Azure SQL Database & Synapse Analytics",
                "Azure Cognitive Search (semantic indexing)"
            ],
            "🔐 Azure Security & DevOps",
            [
                "Azure Active Directory (authentication)",
                "Azure Key Vault (secrets management)",
                "Azure Monitor & Application Insights",
                "Azure DevOps (CI/CD pipelines)",
                "Azure API Management (gateway)",
                "Azure Foundry (development platform)"
            ],
            1.8
        )
        
        # Slide 5: Azure Multi-Tenant Security
        slide = self.add_content_slide(prs, "Azure-Native Multi-Tenant Security Model", 5)
        self.add_icon_bullets(slide, [
            "🔐 Azure Active Directory Integration",
            "  Azure AD B2C for customer identity management",
            "  Multi-tenant Azure AD application registration",
            "  Role-Based Access Control (RBAC) with Azure AD groups",
            "  Conditional access policies and MFA enforcement",
            "🛡️ Azure Security Services",
            "  Azure Key Vault for secrets and certificate management",
            "  Azure SQL Row-Level Security (RLS) policies",
            "  Azure Monitor & Security Center for threat detection",
            "  Azure Policy for compliance enforcement",
            "⚡ Azure Performance & Scale",
            "  Azure Redis Cache for tenant-specific caching",
            "  Azure Load Balancer for traffic distribution",
            "  Azure CDN for global content delivery",
            "  Azure Auto-scaling for elastic compute"
        ], 0.8, 1.8)
        
        # Slide 6: Azure Foundry Development Workflow
        slide = self.add_content_slide(prs, "Azure Foundry Development Pipeline", 6)
        self.add_premium_table(
            slide,
            ['Phase', 'Azure Foundry Feature', 'Azure Services', 'Deliverable'],
            [
                ['Design', 'AI-assisted architecture', 'Azure OpenAI, Cognitive Services', 'Architecture blueprint'],
                ['Development', 'Foundry code generation', 'Azure Repos, GitHub Copilot', 'Application code'],
                ['Testing', 'Automated testing suite', 'Azure Test Plans, Load Testing', 'Quality assurance'],
                ['Deployment', 'CI/CD automation', 'Azure DevOps, Container Registry', 'Production release'],
                ['Monitoring', 'Observability platform', 'Application Insights, Monitor', 'Live telemetry']
            ],
            1.8
        )
        
        # Slide 7: Section - Power BI Integration
        self.add_section_divider(
            prs,
            "Power BI Integration Strategy",
            "Native Microsoft ecosystem connectivity",
            "02"
        )
        
        # Slide 8: Power BI Integration Approaches
        slide = self.add_content_slide(prs, "Azure-Native Power BI Integration", 8)
        self.add_two_column_premium(
            slide,
            "🔌 Custom Power BI Connector",
            [
                "Power Query M language with Azure auth",
                "Native Power BI Desktop & Service experience",
                "Direct Query & Import mode support",
                "Azure AD SSO authentication",
                "Microsoft AppSource certification",
                "Scheduled refresh with Azure credentials"
            ],
            "� Azure API Management Gateway",
            [
                "Azure APIM with OAuth 2.0 authentication",
                "OData/JSON endpoints for Power BI",
                "Rate limiting & throttling policies",
                "API versioning & lifecycle management",
                "Azure Monitor integration for analytics",
                "Power BI Web connector compatibility"
            ],
            1.8
        )
        
        # Slide 9: Azure OpenAI Capabilities
        slide = self.add_content_slide(prs, "Azure OpenAI Service Integration", 9)
        self.add_premium_metrics(
            slide,
            [
                {'value': 'GPT-4o', 'label': 'Azure OpenAI\nModel'},
                {'value': '95%+', 'label': 'SQL Query\nAccuracy'},
                {'value': '99.9%', 'label': 'Azure SLA\nUptime'},
                {'value': '<2s', 'label': 'Response\nTime'}
            ],
            2.2
        )
        self.add_icon_bullets(slide, [
            "🤖 Azure OpenAI Service: Enterprise-grade GPT-4o deployment in Azure",
            "🔒 Data Privacy: Your data never leaves Azure, full compliance with regulations",
            "🌐 Global Availability: Azure regions worldwide for low-latency access",
            "📊 Semantic Understanding: Context-aware SQL generation from natural language",
            "�️ Content Filtering: Azure content safety filters for responsible AI"
        ], 0.8, 5.2)
        
        # Slide 10: Azure Compliance & Governance
        slide = self.add_content_slide(prs, "Azure Compliance & Governance Framework", 10)
        self.add_premium_table(
            slide,
            ['Compliance Area', 'Azure Service', 'Certification', 'Coverage'],
            [
                ['Data Residency', 'Azure Regions', 'GDPR, CCPA', 'Data sovereignty'],
                ['Identity Security', 'Azure AD', 'ISO 27001, SOC 2', 'Authentication'],
                ['Data Encryption', 'Azure Key Vault', 'FIPS 140-2', 'At-rest & in-transit'],
                ['Audit & Logging', 'Azure Monitor', 'HIPAA, PCI-DSS', 'Compliance tracking'],
                ['Threat Protection', 'Azure Security Center', 'CIS Benchmarks', 'Threat detection']
            ],
            1.8
        )
        
        # Slide 11: Section - Microsoft Partnership
        self.add_section_divider(
            prs,
            "Microsoft Partnership Benefits",
            "Leveraging Microsoft ecosystem for competitive advantage",
            "03"
        )
        
        # Slide 12: Azure Partner Advantages
        slide = self.add_content_slide(prs, "Microsoft Azure Partner Benefits", 12)
        self.add_icon_bullets(slide, [
            "🤝 Microsoft Partnership Tier",
            "  Azure Partner designation with Microsoft",
            "  Access to Microsoft co-sell programs",
            "  Azure Marketplace listing eligibility",
            "  Microsoft Field Engineering support",
            "☁️ Azure Technology Benefits",
            "  Azure credits for development & testing",
            "  Priority access to new Azure services",
            "  Azure Foundry platform for rapid development",
            "  Azure OpenAI Service enterprise deployment",
            "� Power Platform Integration",
            "  Native Power BI connector development",
            "  Power Apps & Power Automate integration",
            "  Microsoft AppSource certification path",
            "  Microsoft 365 ecosystem connectivity"
        ], 0.8, 1.8)
        
        # Slide 13: Azure Competitive Advantages
        slide = self.add_content_slide(prs, "Azure-Native Competitive Differentiation", 13)
        self.add_two_column_premium(
            slide,
            "☁️ Azure-First Architecture",
            [
                "100% Azure-native deployment",
                "Built on Azure Foundry platform",
                "Azure OpenAI Service (GPT-4o)",
                "Seamless Azure AD integration",
                "Native Power BI connectivity",
                "Azure Marketplace ready"
            ],
            "🏆 vs. Non-Azure Solutions",
            [
                "No cross-cloud complexity or latency",
                "Single Microsoft billing & support",
                "Microsoft co-sell opportunity eligible",
                "Enterprise Azure SLA guarantees",
                "Azure compliance certifications",
                "Unified Microsoft ecosystem"
            ],
            1.8
        )
        
        # Slide 14: Azure Foundry Development Roadmap
        slide = self.add_content_slide(prs, "Strategic Roadmap with Azure Foundry", 14)
        self.add_icon_bullets(slide, [
            "☁️ Phase 1: Azure Infrastructure (Q1 2025)",
            "  Deploy core services on Azure App Service & AKS",
            "  Configure Azure AD B2C for multi-tenant authentication",
            "  Establish Azure DevOps CI/CD pipelines with Foundry",
            "🔌 Phase 2: Power BI Integration (Q1-Q2 2025)",
            "  Develop custom Power BI connector with Azure auth",
            "  Implement Azure API Management gateway",
            "  Microsoft AppSource certification submission",
            "🚀 Phase 3: Azure Marketplace Launch (Q2 2025)",
            "  Azure Marketplace listing and GTM strategy",
            "  Microsoft co-sell program enrollment",
            "  Enterprise customer pilot program (3-5 customers)",
            "� Phase 4: Scale & Optimize (Q3 2025+)",
            "  Azure global region expansion",
            "  Advanced Azure Monitor & Application Insights",
            "  Azure Cognitive Services integration (speech, vision)"
        ], 0.8, 1.8)
        
        # Slide 15: Closing with Azure branding
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Gradient background
        self.add_gradient_background(slide, self.colors['primary_dark'], self.colors['primary'])
        
        # Decorative elements
        self.add_decorative_elements(slide)
        
        # Large text
        closing_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.2), Inches(9.333), Inches(2.2)
        )
        closing_frame = closing_box.text_frame
        closing_frame.text = "Built on Azure.\nPowered by AI.\nReady for Enterprise."
        closing_frame.word_wrap = True
        
        for para in closing_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(52)
            para.font.bold = True
            para.font.color.rgb = self.colors['white']
            para.font.name = 'Segoe UI'
            para.line_spacing = 1.3
        
        # Accent line
        line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(5), Inches(4.8), Inches(3.333), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent_cyan']
        line.line.fill.background()
        
        # Microsoft Partner badge
        badge_box = slide.shapes.add_textbox(
            Inches(2), Inches(5.4), Inches(9.333), Inches(1.2)
        )
        badge_frame = badge_box.text_frame
        badge_frame.text = "Microsoft Azure Partner\nBuilt on Azure Foundry • Certified for Power BI"
        
        for para in badge_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(20)
            para.font.color.rgb = self.colors['accent_cyan']
            para.font.name = 'Segoe UI Semibold'
        
        # Contact
        contact_box = slide.shapes.add_textbox(
            Inches(2), Inches(6.5), Inches(9.333), Inches(0.6)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "NL2Q Analytics Team • enterprise@nl2q-analytics.com"
        
        p = contact_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['gray_300']
        p.font.name = 'Segoe UI'
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"nl2q_azure_powerbi_{timestamp}.pptx"
        filepath = self.output_dir / filename
        
        prs.save(str(filepath))
        
        print(f"\n✅ Azure-Native Premium PowerPoint generated!")
        print(f"📄 File: {filepath}")
        print(f"📊 Slides: {len(prs.slides)}")
        print(f"☁️ Focus: Azure Foundry, Azure OpenAI, Power BI Integration")
        print(f"🤝 Positioning: Microsoft Azure Partner")
        
        return str(filepath)


def main():
    """Main execution"""
    generator = EnterprisePremiumPPTX()
    output_file = generator.generate_powerbi_deck()
    
    print(f"\n🎉 Azure-Native Premium Presentation Ready!")
    print(f"\n☁️ Azure Features:")
    print("   • Azure Foundry development platform")
    print("   • Azure OpenAI Service (GPT-4o)")
    print("   • Azure Active Directory integration")
    print("   • Azure Key Vault security")
    print("   • Microsoft Partner branding")
    print("   • Power BI native connectivity")
    print("   • No pricing/cost information")
    print(f"\n📂 Location: {output_file}")


if __name__ == "__main__":
    main()
