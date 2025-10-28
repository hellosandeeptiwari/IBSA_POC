"""
Corporate Template PowerPoint Generator
Uses Conexus Corporate Template 2025 for brand consistency

Author: NL2Q Analytics Team
Date: October 2025
Version: 3.0 - Corporate Branding
"""

import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Dict, List, Optional


class CorporatePresentationGenerator:
    """
    Generate presentations using Conexus Corporate Template
    Maintains brand consistency with corporate standards
    """
    
    def __init__(self, template_path: str = r"assets\Conexus Corporate Template 2025.pptx",
                 output_dir: str = "outputs"):
        self.template_path = Path(template_path)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")
        
        # Layout mapping based on template analysis
        self.layouts = {
            'title': 0,              # Title Slide
            'blue_bar': 1,           # Blue bar open content
            'left_callout': 2,       # Left Callout with text
            'title_content': 3,      # Title single Content
            'quote': 4,              # Large Quote
            'single_content': 5,     # Title single Content
            'two_column': 6,         # 2 Column Content
            'three_column': 7,       # 3 Column Content
            'title_only': 8,         # Title
            'title_orange': 9,       # Title Orange
            'business_function': 10, # Business Function
            'icon_listing': 11,      # 4 Icon Listing
            'white_bg': 12,          # White background
            'headline_only': 13,     # Headline only
            'section_blue': 14,      # Section Divider Blue
            'section_orange': 15,    # Section Divider Orange
            'closing': 16            # Closing
        }
        
        # Color scheme for text based on layout type
        # Dark backgrounds need light text, light backgrounds need dark text
        self.text_colors = {
            'dark_bg': RGBColor(255, 255, 255),  # White text for dark backgrounds
            'light_bg': RGBColor(30, 41, 59),    # Dark text for light backgrounds
            'blue_bg': RGBColor(255, 255, 255),  # White text for blue backgrounds
        }
        
        # Map layouts to appropriate text colors
        self.layout_text_colors = {
            'title': self.text_colors['dark_bg'],
            'blue_bar': self.text_colors['light_bg'],      # Blue bar but content area is white
            'left_callout': self.text_colors['light_bg'],
            'title_content': self.text_colors['light_bg'],
            'single_content': self.text_colors['light_bg'],
            'two_column': self.text_colors['light_bg'],
            'three_column': self.text_colors['light_bg'],
            'white_bg': self.text_colors['light_bg'],
            'section_blue': self.text_colors['dark_bg'],
            'section_orange': self.text_colors['dark_bg'],
            'closing': self.text_colors['dark_bg'],
        }
        
    def create_presentation(self) -> Presentation:
        """Load corporate template"""
        print(f"📋 Loading corporate template: {self.template_path.name}")
        prs = Presentation(str(self.template_path))
        
        # Remove example slides, keep only master slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        return prs
    
    def add_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        """Add corporate title slide (Layout 0)"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['title']])
        
        # Find and populate placeholders
        for shape in slide.shapes:
            if shape.has_text_frame:
                if 'title' in shape.name.lower():
                    shape.text_frame.text = title
                elif 'subtitle' in shape.name.lower() or 'text' in shape.name.lower():
                    if subtitle:
                        shape.text_frame.text = subtitle
        
        return slide
    
    def add_section_divider(self, prs: Presentation, section_title: str, 
                           use_orange: bool = False):
        """Add section divider (Layout 14 or 15)"""
        layout_key = 'section_orange' if use_orange else 'section_blue'
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts[layout_key]])
        
        # Populate section title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = section_title
                break
        
        return slide
    
    def add_content_slide(self, prs: Presentation, title: str, 
                         layout_type: str = 'blue_bar') -> any:
        """Add content slide with specified layout"""
        layout_idx = self.layouts.get(layout_type, self.layouts['blue_bar'])
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Set title
        for shape in slide.shapes:
            if shape.has_text_frame and 'title' in shape.name.lower():
                shape.text_frame.text = title
                break
        
        return slide
    
    def add_bullet_content(self, slide, content_items: List[str], 
                          placeholder_idx: int = 13, text_color: RGBColor = None):
        """Add bullet points to content placeholder with proper text color"""
        if text_color is None:
            text_color = RGBColor(30, 41, 59)  # Dark slate for readability
            
        try:
            # Find content placeholder (idx 13 for most content slides)
            text_frame = slide.placeholders[placeholder_idx].text_frame
            text_frame.clear()
            
            for i, item in enumerate(content_items):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # Determine indent level
                level = 0
                text = item
                if item.startswith('  '):
                    level = 1
                    text = item.strip()
                
                p.text = text
                p.level = level
                
                # Set text color explicitly
                for run in p.runs:
                    run.font.color.rgb = text_color
                    
        except Exception as e:
            print(f"⚠️  Could not add bullets (idx {placeholder_idx}): {e}")
            # Try alternative approach with text boxes
            try:
                for shape in slide.shapes:
                    if shape.has_text_frame and 'Text Placeholder' in shape.name:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        for i, item in enumerate(content_items):
                            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                            p.text = item.strip() if item.startswith('  ') else item
                            p.level = 1 if item.startswith('  ') else 0
                            # Set text color
                            for run in p.runs:
                                run.font.color.rgb = text_color
                        break
            except Exception as e2:
                print(f"⚠️  Fallback also failed: {e2}")
    
    def add_two_column_content(self, prs: Presentation, title: str,
                              left_items: List[str], right_items: List[str],
                              text_color: RGBColor = None):
        """Add two-column layout slide with proper text color"""
        if text_color is None:
            text_color = RGBColor(30, 41, 59)  # Dark slate for readability
            
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['two_column']])
        
        # Set title (placeholder 0)
        slide.placeholders[0].text_frame.text = title
        
        # Add content to left and right columns
        # Layout 6 has placeholders: 0=Title, 1=Left Content, 13=Text, 14=Right Content
        try:
            # Left column (placeholder 1)
            left_frame = slide.placeholders[1].text_frame
            left_frame.clear()
            for i, item in enumerate(left_items):
                p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
                if item.startswith('  '):
                    p.level = 1
                    p.text = item.strip()
                else:
                    p.level = 0
                    p.text = item
                # Set text color
                for run in p.runs:
                    run.font.color.rgb = text_color
            
            # Right column (placeholder 14)
            right_frame = slide.placeholders[14].text_frame
            right_frame.clear()
            for i, item in enumerate(right_items):
                p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
                if item.startswith('  '):
                    p.level = 1
                    p.text = item.strip()
                else:
                    p.level = 0
                    p.text = item
                # Set text color
                for run in p.runs:
                    run.font.color.rgb = text_color
        except Exception as e:
            print(f"⚠️  Could not populate columns: {e}")
            import traceback
            traceback.print_exc()
        
        return slide
    
    def add_closing_slide(self, prs: Presentation, closing_text: str = ""):
        """Add corporate closing slide (Layout 16)"""
        slide = prs.slides.add_slide(prs.slide_layouts[self.layouts['closing']])
        
        if closing_text:
            for shape in slide.shapes:
                if shape.has_text_frame and not shape.is_placeholder:
                    shape.text_frame.text = closing_text
                    break
        
        return slide
    
    def add_custom_text_shape(self, slide, text: str, left: float, top: float,
                             width: float, height: float, font_size: int = 18,
                             bold: bool = False, color: RGBColor = None):
        """Add custom text box to slide"""
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        
        return text_box
    
    def generate_azure_powerbi_deck(self) -> str:
        """Generate complete Azure Power BI Integration deck using corporate template"""
        print("🎨 Generating Corporate Branded PowerPoint...")
        print("🏢 Using Conexus Corporate Template 2025")
        
        prs = self.create_presentation()
        
        # Slide 1: Corporate Title
        print("  📄 Adding title slide...")
        self.add_title_slide(
            prs,
            "Azure Power BI Integration Strategy",
            "NL2Q Analyst Platform | Built on Azure AI Foundry"
        )
        
        # Slide 2: Executive Summary
        print("  📄 Adding executive summary...")
        slide = self.add_content_slide(prs, "Executive Summary", 'blue_bar')
        self.add_bullet_content(slide, [
            "🚀 Azure-Native Integration: Seamless Power BI connectivity with NL2Q Analyst",
            "🔧 Three Implementation Approaches: Custom Connector, REST API, Embedded Analytics",
            "☁️ Azure AI Foundry: Built on Microsoft's premier AI development platform",
            "🔒 Enterprise Security: Azure AD, Key Vault, and multi-tenant isolation",
            "⚡ Enhanced Productivity: 40% reduction in report creation cycles",
            "🌐 Microsoft Partnership: Leveraging Azure ecosystem advantages"
        ], text_color=self.layout_text_colors['blue_bar'])
        
        # Slide 3: Section - Architecture
        print("  📄 Adding architecture section...")
        self.add_section_divider(prs, "Azure-Native Technical Architecture")
        
        # Slide 4: Integration Approaches
        print("  📄 Adding integration approaches...")
        self.add_two_column_content(
            prs,
            "Three Strategic Integration Approaches",
            [
                "🔌 Approach 1: Custom Connector",
                "Power Query M language",
                "Native Power BI experience",
                "Azure AD authentication",
                "Certified connector ready",
                "Seamless data refresh"
            ],
            [
                "🌐 Approach 2: Azure API Management",
                "OData/JSON endpoints",
                "API Gateway integration",
                "Rate limiting & throttling",
                "Developer portal access",
                "Multi-platform support"
            ],
            text_color=self.layout_text_colors['two_column']
        )
        
        # Slide 5: Azure AI Foundry Architecture
        print("  📄 Adding Azure AI Foundry details...")
        slide = self.add_content_slide(prs, "Built on Azure AI Foundry", 'blue_bar')
        self.add_bullet_content(slide, [
            "🤖 Azure OpenAI Service: GPT-4o and GPT-4o-mini for SQL generation",
            "  Prompt flow orchestration for complex query workflows",
            "  Responsible AI guardrails and content filtering",
            "🔍 Azure AI Search: Vector search for semantic schema understanding",
            "  Hybrid search combining keyword and semantic",
            "  Integrated with Azure OpenAI embeddings",
            "📊 Azure Machine Learning: Model monitoring and evaluation",
            "  Performance tracking and accuracy metrics",
            "  Continuous improvement pipeline",
            "🔐 Azure Security: End-to-end protection and compliance",
            "  Azure AD for authentication and authorization",
            "  Azure Key Vault for credential management",
            "  Azure Monitor for observability and alerting"
        ], text_color=self.layout_text_colors['blue_bar'])
        
        # Slide 6: Multi-Tenant Security
        print("  📄 Adding multi-tenant architecture...")
        slide = self.add_content_slide(prs, "Azure Multi-Tenant Security Model", 'blue_bar')
        self.add_bullet_content(slide, [
            "🔐 Three-Tier Azure Isolation Strategy",
            "  Tier 1: Snowflake schema-per-tenant (data isolation)",
            "  Tier 2: Azure SQL Row-Level Security (access control)",
            "  Tier 3: Azure AD B2C with tenant context",
            "☁️ Azure Security Services",
            "  Azure Key Vault: Secure credential storage",
            "  Azure AD: Role-Based Access Control (RBAC)",
            "  Azure Monitor: Audit logging and compliance",
            "  Azure Sentinel: Threat detection and response",
            "⚡ Azure Performance Services",
            "  Azure Cache for Redis: Per-tenant caching",
            "  Azure Front Door: Global load balancing",
            "  Azure CDN: Content delivery optimization"
        ], text_color=self.layout_text_colors['blue_bar'])
        
        # Slide 7: Implementation Timeline
        print("  📄 Adding implementation timeline...")
        slide = self.add_content_slide(prs, "Azure Implementation Roadmap", 'blue_bar')
        self.add_bullet_content(slide, [
            "Phase 1: Azure Foundation (4 weeks)",
            "  Azure AD B2C configuration and tenant setup",
            "  Azure Key Vault integration for secrets",
            "  Azure Monitor and Application Insights setup",
            "Phase 2: Azure API Development (3 weeks)",
            "  Azure API Management gateway configuration",
            "  Azure Functions for serverless processing",
            "  Power BI Custom Connector development",
            "Phase 3: Azure AI Integration (3 weeks)",
            "  Azure OpenAI Service prompt optimization",
            "  Azure AI Search index configuration",
            "  User acceptance testing on Azure",
            "Phase 4: Azure Deployment (2 weeks)",
            "  Azure DevOps CI/CD pipeline setup",
            "  Production deployment to Azure regions",
            "  Azure monitoring and alerting configuration"
        ], text_color=self.layout_text_colors['blue_bar'])
        
        # Slide 8: Section - Azure Capabilities
        print("  📄 Adding capabilities section...")
        self.add_section_divider(prs, "Azure Technical Capabilities", use_orange=True)
        
        # Slide 9: Azure Stack
        print("  📄 Adding Azure stack details...")
        self.add_two_column_content(
            prs,
            "Complete Azure Technology Stack",
            [
                "☁️ Azure Compute & Hosting",
                "Azure App Service (FastAPI backend)",
                "Azure Static Web Apps (React frontend)",
                "Azure Functions (serverless processing)",
                "Azure Container Instances (scaling)",
                "",
                "🤖 Azure AI & Data Services",
                "Azure OpenAI Service (GPT-4o/o3-mini)",
                "Azure AI Search (vector search)",
                "Azure Machine Learning (monitoring)",
                "Azure Cognitive Services"
            ],
            [
                "🔐 Azure Security & Identity",
                "Azure AD B2C (authentication)",
                "Azure Key Vault (secrets management)",
                "Azure Private Link (network security)",
                "Azure Security Center (compliance)",
                "",
                "📊 Azure Data & Integration",
                "Azure API Management (gateway)",
                "Azure Cache for Redis (performance)",
                "Azure Monitor (observability)",
                "Azure DevOps (CI/CD)"
            ]
        )
        
        # Slide 10: Feasibility Analysis
        print("  📄 Adding feasibility analysis...")
        self.add_two_column_content(
            prs,
            "Azure Technical Readiness Assessment",
            [
                "✅ Azure Strengths",
                "Production FastAPI on Azure App Service",
                "Azure OpenAI proven accuracy (GPT-4o)",
                "Snowflake on Azure integration",
                "React frontend on Azure Static Web Apps",
                "Azure AI Search optimized",
                "Azure DevOps pipelines configured"
            ],
            [
                "🎯 Azure Enhancement Areas",
                "Azure AD B2C implementation",
                "Azure API Management setup",
                "Azure Monitor comprehensive logging",
                "Azure Key Vault full integration",
                "Power BI Custom Connector",
                "Azure Load Testing validation"
            ]
        )
        
        # Slide 11: Performance Metrics
        print("  📄 Adding performance metrics...")
        slide = self.add_content_slide(prs, "Azure Performance Indicators", 'blue_bar')
        self.add_bullet_content(slide, [
            "📈 Performance Targets",
            "  12-16 weeks to Azure production deployment",
            "  40% efficiency improvement in report creation",
            "  99.9% uptime SLA with Azure availability zones",
            "  <2 second query response with Azure Cache for Redis",
            "🎯 Azure Advantages",
            "  Global scale with Azure regions worldwide",
            "  Enterprise SLAs and support",
            "  Integrated security and compliance",
            "  Cost optimization with Azure Reserved Instances",
            "💡 Microsoft Partnership Benefits",
            "  Priority Azure technical support",
            "  Co-sell opportunities in Azure Marketplace",
            "  Access to Microsoft field resources",
            "  Azure credits and incentives"
        ], text_color=self.layout_text_colors['blue_bar'])
        
        # Slide 12: Section - Business Value
        print("  📄 Adding business value section...")
        self.add_section_divider(prs, "Business Impact & Azure Value")
        
        # Slide 13: Business Benefits
        print("  📄 Adding business benefits...")
        slide = self.add_content_slide(prs, "Expected Business Outcomes", 'blue_bar')
        self.add_bullet_content(slide, [
            "📈 Operational Excellence",
            "  40% reduction in report creation time",
            "  60% faster ad-hoc query responses",
            "  70% decrease in SQL developer dependency",
            "  Azure-scale performance and reliability",
            "🎯 Enhanced User Experience",
            "  Natural language queries in Power BI",
            "  Self-service analytics for business users",
            "  Azure AD single sign-on integration",
            "  Global availability via Azure regions",
            "💼 Strategic Azure Advantages",
            "  Microsoft partnership benefits",
            "  Azure Marketplace presence",
            "  Enterprise support and SLAs",
            "  Seamless Microsoft 365 integration"
        ], text_color=self.layout_text_colors['blue_bar'])
        
        # Slide 14: Competitive Positioning
        print("  📄 Adding competitive positioning...")
        self.add_two_column_content(
            prs,
            "Azure Partnership Competitive Edge",
            [
                "🚀 Azure-Native Capabilities",
                "Built on Azure AI Foundry platform",
                "Azure OpenAI GPT-4o integration",
                "Azure AI Search semantic understanding",
                "Azure AD enterprise authentication",
                "Native Power BI connectivity",
                "Multi-cloud Azure, Snowflake support"
            ],
            [
                "🏆 Microsoft Partnership Advantages",
                "Azure Marketplace listing",
                "Co-sell with Microsoft field",
                "Priority technical support",
                "Access to Azure credits",
                "Joint marketing opportunities",
                "Enterprise customer access"
            ],
            text_color=self.layout_text_colors['two_column']
        )
        
        # Slide 15: Strategic Recommendations
        print("  📄 Adding recommendations...")
        slide = self.add_content_slide(prs, "Strategic Recommendations", 'blue_bar')
        self.add_bullet_content(slide, [
            "🎯 Phase 1: Azure API Management Gateway (fastest time-to-value)",
            "📅 Timeline: Q1 2025 start, Q2 2025 Azure production-ready",
            "💼 Resources: 3 Azure-certified engineers, 1 Azure DevOps specialist",
            "🔐 Security First: Azure AD B2C and Key Vault before integration",
            "👥 Pilot: 3 enterprise customers on Azure for beta validation",
            "📊 Success Metrics: Query accuracy, Azure performance, user adoption",
            "🚀 Azure Marketplace: List as Power BI certified solution",
            "📚 Enablement: Azure architecture documentation and runbooks"
        ], text_color=self.layout_text_colors['blue_bar'])
        
        # Slide 16: Closing
        print("  📄 Adding closing slide...")
        self.add_closing_slide(prs)
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"nl2q_azure_powerbi_{timestamp}.pptx"
        filepath = self.output_dir / filename
        
        prs.save(str(filepath))
        
        print(f"\n✅ Corporate branded presentation generated!")
        print(f"📄 File: {filepath}")
        print(f"📊 Slides: {len(prs.slides)}")
        print(f"🏢 Template: Conexus Corporate Template 2025")
        print(f"☁️  Focus: Azure AI Foundry & Microsoft Partnership")
        
        return str(filepath)


def main():
    """Main execution"""
    try:
        generator = CorporatePresentationGenerator()
        output_file = generator.generate_azure_powerbi_deck()
        
        print(f"\n🎉 Success! Your corporate-branded presentation is ready.")
        print(f"\n💼 Features:")
        print("   ✅ Conexus corporate branding")
        print("   ✅ Azure AI Foundry focus")
        print("   ✅ Microsoft partnership emphasis")
        print("   ✅ Azure-native architecture")
        print("   ✅ No pricing information")
        print(f"\n📂 Location: {output_file}")
        
    except FileNotFoundError as e:
        print(f"\n❌ Error: {e}")
        print("💡 Make sure 'Conexus Corporate Template 2025.pptx' is in the assets folder")
    except Exception as e:
        print(f"\n❌ Unexpected error: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
