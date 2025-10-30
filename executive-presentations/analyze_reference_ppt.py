"""Analyze reference presentation to see what EDA elements they used"""
from pptx import Presentation
from pathlib import Path

ref_path = Path("assets/Quarterly Business Summary 7.2025 v3-forConexus 1.pptx")
prs = Presentation(ref_path)

print(f"\n{'='*60}")
print(f"REFERENCE PRESENTATION ANALYSIS")
print(f"{'='*60}")
print(f"Total slides: {len(prs.slides)}\n")

print("=== SLIDE TITLES ===")
for i, slide in enumerate(prs.slides, 1):
    if slide.shapes.title:
        print(f"{i:2d}. {slide.shapes.title.text}")
    else:
        print(f"{i:2d}. (No title)")

print("\n=== SLIDES WITH CHARTS/IMAGES ===")
for i, slide in enumerate(prs.slides, 1):
    title = slide.shapes.title.text if slide.shapes.title else "(No title)"
    
    # Count different types of content
    charts = 0
    images = 0
    tables = 0
    textboxes = 0
    
    for shape in slide.shapes:
        if shape.shape_type == 3:  # Chart
            charts += 1
        elif shape.shape_type == 13:  # Picture
            images += 1
        elif shape.shape_type == 19:  # Table
            tables += 1
        elif shape.shape_type == 17:  # Textbox
            textboxes += 1
    
    if charts > 0 or images > 0 or tables > 0:
        print(f"\n{i:2d}. {title}")
        if charts > 0:
            print(f"    📊 Charts: {charts}")
        if images > 0:
            print(f"    🖼️  Images: {images}")
        if tables > 0:
            print(f"    📋 Tables: {tables}")
        print(f"    📝 Text boxes: {textboxes}")

print("\n" + "="*60)
